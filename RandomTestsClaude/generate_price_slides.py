import sys, copy
sys.stdout.reconfigure(encoding='utf-8')

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ─── PATHS ──────────────────────────────────────────────────────────────
EXCEL_PATH    = r'D:\RandomTestsClaude\North TMT pricing_9th Mar.xlsx'
TEMPLATE_PATH = r'D:\RandomTestsClaude\Haryana_NewSlides.pptx'
OUTPUT_PATH   = r'D:\RandomTestsClaude\Delhi_PriceMapping.pptx'
SHEET_NAME    = '1_Summary at PL'

# ─── COLORS (from Haryana design system) ────────────────────────────────
BLUE       = RGBColor(0x18, 0x48, 0x9D)
DARK_NAVY  = RGBColor(0x0D, 0x2B, 0x5E)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TITLE  = RGBColor(0x7F, 0x7F, 0x7F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xD9, 0x77, 0x06)
LIGHT_SLATE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN       = RGBColor(0x05, 0x72, 0x3A)
TABLE_ALT   = RGBColor(0xF0, 0xF4, 0xF8)

# Gap badge colors
GREEN_BADGE  = RGBColor(0x05, 0x72, 0x3A)   # Competitive
AMBER_BADGE  = RGBColor(0xD9, 0x77, 0x06)   # Watch
RED_BADGE    = RGBColor(0xDC, 0x26, 0x26)   # Action Needed

FONT = 'Calibri'

# ─── Layout constants ───────────────────────────────────────────────────
MARGIN_L    = Emu(330000)
FULL_W      = Emu(11530000)
SLIDE_W     = Emu(12192000)

# ─── MARKET CONFIG ──────────────────────────────────────────────────────
MARKET_COLUMNS = {
    'DL': {'name': 'Delhi (DL)',         'distr_col': 2,  'dealer_col': 3,  'gap_col': 4},
    'HR': {'name': 'Haryana (HR)',       'distr_col': 5,  'dealer_col': 6,  'gap_col': 7},
    'RJ': {'name': 'Rajasthan (RJ)',     'distr_col': 8,  'dealer_col': 9,  'gap_col': 10},
    'UP': {'name': 'UP East+Central',    'distr_col': 11, 'dealer_col': 12, 'gap_col': 13},
    'CH': {'name': 'Chandigarh (CH)',    'distr_col': 14, 'dealer_col': 15, 'gap_col': 16},
    'PB': {'name': 'Punjab (PB)',        'distr_col': 17, 'dealer_col': 18, 'gap_col': 19},
    'UK': {'name': 'Uttarakhand (UK)',   'distr_col': 20, 'dealer_col': 21, 'gap_col': 22},
    'HP': {'name': 'Himachal Pradesh',   'distr_col': 23, 'dealer_col': 24, 'gap_col': 25},
    'JK': {'name': 'Jammu & Kashmir',    'distr_col': 26, 'dealer_col': 27, 'gap_col': 28},
}

BRAND_ROWS = range(4, 13)  # rows 4 through 12
SKIP_ROWS  = {9}           # <General SAIL position> note row


# ─── HELPER FUNCTIONS (from Haryana design system) ──────────────────────

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_title(slide, text, subtitle=None):
    """Gray 20pt bold title at top of slide, with optional right-aligned subtitle."""
    txb = slide.shapes.add_textbox(Emu(160000), Emu(120000), Emu(8000000), Emu(380000))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = GRAY_TITLE
    if subtitle:
        txb2 = slide.shapes.add_textbox(Emu(8500000), Emu(180000), Emu(3500000), Emu(300000))
        tf2 = txb2.text_frame
        tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
        r2 = tf2.paragraphs[0].add_run()
        r2.text = subtitle
        r2.font.name = FONT; r2.font.size = Pt(12); r2.font.bold = False; r2.font.color.rgb = DARK_TEXT

def add_slide_num(slide, n):
    txb = slide.shapes.add_textbox(Emu(9448800), Emu(6484939), Emu(2743200), Emu(365125))
    tf = txb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run()
    r.text = str(n)
    r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = GRAY_TITLE

def add_navy_bar(slide, text, y=Emu(6050000)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Emu(420000))
    set_fill(bar, DARK_NAVY); no_line(bar)
    tf = bar.text_frame
    tf.margin_left = Emu(400000); tf.margin_top = Emu(50000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = WHITE

def add_section_banner(slide, x, y, w, text, fill=BLUE, h=Emu(270000)):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(rect, fill); no_line(rect)
    rect.adjustments[0] = 0.06
    tf = rect.text_frame
    tf.margin_left = Emu(100000); tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE

def set_cell_bg(cell, color):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': str(color)}))
    tcPr.append(sf)

def style_cell(cell, text, bold=False, color=DARK_TEXT, bg=None, size=Pt(10), align=None):
    cell.text = ''
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(45720); tf.margin_right = Emu(45720)
    tf.margin_top = Emu(18288); tf.margin_bottom = Emu(18288)
    if align:
        tf.paragraphs[0].alignment = align
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.name = FONT; r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    if bg: set_cell_bg(cell, bg)


# ─── GAP BADGE LOGIC ───────────────────────────────────────────────────

def format_gap_badge(gap):
    """Return (display_string, RGBColor, status_label) for the gap value."""
    if gap is None or gap == 0:
        return ("Benchmark", GREEN_BADGE, "BENCHMARK")
    abs_gap = abs(gap)
    gap_str = f"{int(gap):+,}"
    if abs_gap <= 2000:
        return (f"{gap_str}  Competitive", GREEN_BADGE, "COMPETITIVE")
    elif abs_gap <= 4000:
        return (f"{gap_str}  Watch", AMBER_BADGE, "WATCH")
    else:
        return (f"{gap_str}  Action Needed", RED_BADGE, "ACTION NEEDED")


# ─── DATA READER ────────────────────────────────────────────────────────

def read_market_data(ws, market_key):
    """Read brand pricing data for a single market from the Excel sheet."""
    cfg = MARKET_COLUMNS[market_key]
    brands = []
    for row in BRAND_ROWS:
        if row in SKIP_ROWS:
            continue
        brand_name = ws.cell(row=row, column=1).value
        if not brand_name or str(brand_name).startswith('<'):
            continue
        distr = ws.cell(row=row, column=cfg['distr_col']).value
        dealer = ws.cell(row=row, column=cfg['dealer_col']).value
        gap   = ws.cell(row=row, column=cfg['gap_col']).value
        # Handle 'NA' or non-numeric values
        if isinstance(distr, str): distr = None
        if isinstance(dealer, str): dealer = None
        if isinstance(gap, str): gap = None
        brands.append({
            'name':  str(brand_name).strip(),
            'distr': distr,
            'dealer': dealer,
            'gap':   gap if gap is not None else 0,
        })
    return brands


def get_date_string(ws):
    """Extract date string from cell A3."""
    val = ws.cell(row=3, column=1).value or ''
    return str(val).strip('<>').strip()


# ─── SLIDE BUILDER ──────────────────────────────────────────────────────

def add_pricing_table(slide, data):
    """Add brand comparison table at top of slide."""
    TABLE_X = MARGIN_L
    TABLE_Y = Emu(940000)
    TABLE_W = FULL_W
    num_rows = len(data) + 1
    ROW_H = Emu(225000)
    TABLE_H = ROW_H * num_rows

    COL_WIDTHS = [Emu(3000000), Emu(2500000), Emu(2500000), Emu(3530000)]

    ts = slide.shapes.add_table(num_rows, 4, TABLE_X, TABLE_Y, TABLE_W, TABLE_H)
    tbl = ts.table

    for i, w in enumerate(COL_WIDTHS):
        tbl.columns[i].width = w

    # Header row
    headers = ['Brand', 'Distributor (Rs)', 'Dealer (Rs)', 'Gap vs TISCON']
    for ci, h in enumerate(headers):
        style_cell(tbl.cell(0, ci), h, bold=True, color=WHITE, bg=BLUE, size=Pt(10))

    # Data rows
    for ri, brand in enumerate(data):
        row_idx = ri + 1
        bg = TABLE_ALT if row_idx % 2 == 0 else None
        is_jsw = 'JSW' in brand['name']

        # Brand name
        style_cell(tbl.cell(row_idx, 0), brand['name'],
                   bold=is_jsw, color=DARK_TEXT, bg=bg, size=Pt(10))

        # Distributor price
        distr_str = f"{brand['distr']:,.0f}" if brand['distr'] else 'N/A'
        style_cell(tbl.cell(row_idx, 1), distr_str,
                   color=DARK_TEXT, bg=bg, size=Pt(10), align=PP_ALIGN.RIGHT)

        # Dealer price
        dealer_str = f"{brand['dealer']:,.0f}" if brand['dealer'] else 'N/A'
        style_cell(tbl.cell(row_idx, 2), dealer_str,
                   color=DARK_TEXT, bg=bg, size=Pt(10), align=PP_ALIGN.RIGHT)

        # Gap with color coding
        gap_text, gap_color, _ = format_gap_badge(brand['gap'])
        style_cell(tbl.cell(row_idx, 3), gap_text,
                   bold=True, color=gap_color, bg=bg, size=Pt(10))


def add_price_chart(slide, data):
    """Add grouped bar chart with TISCON reference line."""
    CHART_X = MARGIN_L
    CHART_Y = Emu(3370000)
    CHART_W = Emu(7200000)
    CHART_H = Emu(2550000)

    chart_data = CategoryChartData()
    brand_names = [b['name'] for b in data]
    chart_data.categories = brand_names

    distr_values = tuple(b['distr'] or 0 for b in data)
    dealer_values = tuple(b['dealer'] or 0 for b in data)
    # TISCON dealer price as reference line
    tiscon_dealer = data[0]['dealer'] if data[0]['dealer'] else 0
    tiscon_ref = tuple([tiscon_dealer] * len(data))

    chart_data.add_series('Distributor', distr_values)
    chart_data.add_series('Dealer', dealer_values)
    chart_data.add_series('TISCON Ref', tiscon_ref)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        CHART_X, CHART_Y, CHART_W, CHART_H,
        chart_data
    )
    chart = chart_frame.chart

    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = FONT
    chart.legend.font.size = Pt(8)

    plot = chart.plots[0]
    plot.gap_width = 80

    # Series 0 - Distributor: Blue
    s0 = plot.series[0]
    s0.format.fill.solid()
    s0.format.fill.fore_color.rgb = BLUE

    # Series 1 - Dealer: Amber/Orange
    s1 = plot.series[1]
    s1.format.fill.solid()
    s1.format.fill.fore_color.rgb = ORANGE

    # Value axis
    value_axis = chart.value_axis
    # Set minimum to something reasonable (below lowest value)
    min_val = min(b['distr'] or 99999 for b in data)
    min_val = min(min_val, min(b['dealer'] or 99999 for b in data))
    value_axis.minimum_scale = max(0, int(min_val - 2000))
    value_axis.has_title = False
    value_axis.tick_labels.font.name = FONT
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.number_format = '#,##0'

    # Category axis
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.name = FONT
    cat_axis.tick_labels.font.size = Pt(7)

    # Convert 3rd series (TISCON Ref) to line overlay via XML
    _convert_series_to_line(chart)


def _convert_series_to_line(chart):
    """Move the 3rd bar series to a line chart overlay via XML manipulation."""
    chart_space = chart._chartSpace
    plot_area = chart_space.find(qn('c:chart')).find(qn('c:plotArea'))
    bar_chart = plot_area.find(qn('c:barChart'))

    if bar_chart is None:
        return

    series_list = bar_chart.findall(qn('c:ser'))
    if len(series_list) < 3:
        return

    target_ser = series_list[2]  # 3rd series (TISCON Ref)
    bar_chart.remove(target_ser)

    # Create lineChart element
    line_chart = bar_chart.makeelement(qn('c:lineChart'), {})
    grouping = line_chart.makeelement(qn('c:grouping'), {'val': 'standard'})
    line_chart.append(grouping)
    line_chart.append(target_ser)

    # Set marker to none
    marker = target_ser.makeelement(qn('c:marker'), {})
    symbol = marker.makeelement(qn('c:symbol'), {'val': 'none'})
    marker.append(symbol)
    target_ser.append(marker)

    # Style the line: red dashed
    spPr = target_ser.find(qn('c:spPr'))
    if spPr is None:
        spPr = target_ser.makeelement(qn('c:spPr'), {})
        target_ser.append(spPr)

    ln = spPr.makeelement(qn('a:ln'), {'w': '19050'})
    sf = ln.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:srgbClr'), {'val': 'E8272A'})  # JSW Red
    sf.append(clr)
    ln.append(sf)

    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)
    spPr.append(ln)

    # No fill on bars (since it's now a line)
    noFill = spPr.makeelement(qn('a:noFill'), {})
    spPr.insert(0, noFill)

    # Insert lineChart after barChart
    bar_chart.addnext(line_chart)

    # Copy axis references so the line shares the same axes
    ax_ids = bar_chart.findall(qn('c:axId'))
    for ax in ax_ids:
        line_chart.append(copy.deepcopy(ax))


def add_jsw_spotlight(slide, data):
    """Add JSW Brands Spotlight boxes to the right of the chart."""
    SPOT_X = Emu(7730000)
    SPOT_W = Emu(4130000)

    # Spotlight sub-banner
    add_section_banner(slide, SPOT_X, Emu(3370000), SPOT_W,
                       "JSW Brands Spotlight", fill=DARK_NAVY, h=Emu(270000))

    # Find JSW brands
    jsw_brands = [b for b in data if 'JSW' in b['name']]

    BOX_Y_START = Emu(3700000)
    BOX_H = Emu(1050000)
    BOX_GAP = Emu(100000)

    for i, brand in enumerate(jsw_brands):
        box_y = BOX_Y_START + i * (BOX_H + BOX_GAP)
        gap = brand['gap']
        abs_gap = abs(gap) if gap else 0

        # Determine accent color based on gap severity
        if abs_gap <= 2000:
            accent_color = GREEN
            status = "COMPETITIVE"
        elif abs_gap <= 4000:
            accent_color = ORANGE
            status = "WATCH"
        else:
            accent_color = RED_BADGE
            status = "ACTION NEEDED"

        # Rounded rectangle box
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SPOT_X, box_y, SPOT_W, BOX_H
        )
        set_fill(rect, accent_color); no_line(rect)
        rect.adjustments[0] = 0.04

        tf = rect.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(120000)
        tf.margin_top = Emu(50000)
        tf.margin_right = Emu(80000)

        # Brand name (shortened)
        short_name = brand['name'].replace('JSW One ', '')
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = short_name
        r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE

        # Gap value
        p2 = tf.add_paragraph(); p2.space_before = Pt(6)
        r2 = p2.add_run()
        gap_str = f"{int(gap):+,}" if gap else "0"
        r2.text = f"Gap vs TISCON: {gap_str}"
        r2.font.name = FONT; r2.font.size = Pt(11); r2.font.color.rgb = WHITE

        # Price details
        p3 = tf.add_paragraph(); p3.space_before = Pt(4)
        r3 = p3.add_run()
        distr_str = f"{brand['distr']:,.0f}" if brand['distr'] else 'N/A'
        dealer_str = f"{brand['dealer']:,.0f}" if brand['dealer'] else 'N/A'
        r3.text = f"Distr: {distr_str}  |  Dealer: {dealer_str}"
        r3.font.name = FONT; r3.font.size = Pt(9); r3.font.color.rgb = LIGHT_SLATE

        # Status label
        p4 = tf.add_paragraph(); p4.space_before = Pt(6)
        r4 = p4.add_run()
        r4.text = f"STATUS: {status}"
        r4.font.name = FONT; r4.font.size = Pt(10); r4.font.bold = True; r4.font.color.rgb = WHITE


def generate_footer_text(market_name, data):
    """Auto-generate a data-driven summary for the footer bar."""
    jsw_brands = [b for b in data if 'JSW' in b['name']]
    if not jsw_brands:
        return f"{market_name}: Market pricing overview"

    worst_gap = min(b['gap'] for b in jsw_brands if b['gap'] is not None)
    worst_brand = next(b for b in jsw_brands if b['gap'] == worst_gap)
    short_name = worst_brand['name'].replace('JSW One ', '')

    _, _, status = format_gap_badge(worst_gap)
    return f"{market_name}: {short_name} at {int(worst_gap):+,} vs TISCON  |  {status} — pricing realignment needed"


def build_market_slide(slide, slide_num, market_name, date_str, data):
    """Build a complete pricing dashboard slide for one market."""
    # Title row
    add_title(slide, market_name, subtitle=date_str)
    add_slide_num(slide, slide_num)

    # Section 1: Brand comparison table
    add_section_banner(slide, MARGIN_L, Emu(620000), FULL_W,
                       "Brand Comparison Table")
    add_pricing_table(slide, data)

    # Section 2 banner (full width)
    add_section_banner(slide, MARGIN_L, Emu(3045000), FULL_W,
                       "Price Comparison & JSW Spotlight")

    # Chart (left)
    add_price_chart(slide, data)

    # JSW Spotlight (right)
    add_jsw_spotlight(slide, data)

    # Footer bar
    footer_text = generate_footer_text(market_name, data)
    add_navy_bar(slide, footer_text)


# ─── MAIN ───────────────────────────────────────────────────────────────

def main(markets=None):
    if markets is None:
        markets = ['DL']

    print("Loading Excel data...")
    wb = openpyxl.load_workbook(EXCEL_PATH, data_only=True)
    ws = wb[SHEET_NAME]
    date_str = get_date_string(ws)
    print(f"  Date: {date_str}")

    print("Loading template...")
    prs = Presentation(TEMPLATE_PATH)

    # Find Blank layout
    blank_layout = None
    for l in prs.slide_layouts:
        if l.name == 'Blank':
            blank_layout = l; break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[0]

    # Delete existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    print(f"\nBuilding {len(markets)} slide(s)...\n")

    for i, mkt in enumerate(markets):
        cfg = MARKET_COLUMNS[mkt]
        data = read_market_data(ws, mkt)
        slide = prs.slides.add_slide(blank_layout)
        build_market_slide(slide, i + 1, cfg['name'], date_str, data)
        print(f"  Slide {i+1}: {cfg['name']} done")
        # Print data summary
        for b in data:
            gap_text, _, _ = format_gap_badge(b['gap'])
            print(f"    {b['name']:25s}  Distr: {b['distr'] or 'N/A':>8}  Dealer: {b['dealer'] or 'N/A':>8}  {gap_text}")

    prs.save(OUTPUT_PATH)
    print(f"\nSaved: {OUTPUT_PATH}")
    print(f"Total: {len(prs.slides)} slide(s)")


if __name__ == '__main__':
    main()
