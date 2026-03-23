import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── DESIGN CONSTANTS (from main file analysis) ────────────────────────
BLUE = RGBColor(0x18, 0x48, 0x9D)
DARK_NAVY = RGBColor(0x0D, 0x2B, 0x5E)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TITLE = RGBColor(0x7F, 0x7F, 0x7F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
LIGHT_BLUE = RGBColor(0x93, 0xC5, 0xFD)
TABLE_ALT = RGBColor(0xF0, 0xF4, 0xF8)
GREEN_ACCENT = RGBColor(0x05, 0x72, 0x3A)

FONT = 'Calibri'
SLIDE_W = 12192000
SLIDE_H = 6858000


def set_fill(shape, color):
    """Set solid fill on a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_no_line(shape):
    """Remove outline from shape."""
    shape.line.fill.background()


def add_title(slide, text):
    """Title: Calibri 20pt Bold gray at top."""
    txb = slide.shapes.add_textbox(Emu(160412), Emu(260551), Emu(12057374), Emu(369332))
    tf = txb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = GRAY_TITLE


def add_slide_number(slide, num):
    """Slide number bottom-right."""
    txb = slide.shapes.add_textbox(Emu(9448800), Emu(6484939), Emu(2743200), Emu(365125))
    tf = txb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run()
    r.text = str(num)
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY_TITLE


def add_numbered_circle(slide, left, top, number, size=Emu(268224), fill_color=BLUE):
    """Numbered circle: dark fill, white bold number."""
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    set_fill(oval, fill_color)
    set_no_line(oval)
    tf = oval.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(number)
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_section_banner(slide, left, top, width, height, text, fill_color=BLUE):
    """Blue filled rectangle with white bold text — section divider."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(rect, fill_color)
    set_no_line(rect)
    # Reduce corner rounding
    rect.adjustments[0] = 0.05
    tf = rect.text_frame
    tf.margin_left = Emu(91440)
    tf.margin_top = Emu(27432)
    tf.margin_bottom = Emu(27432)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE


def add_blue_block(slide, left, top, width, height, number, title, bullets):
    """Pattern A: Blue filled block with numbered circle + title + bullet text."""
    # Background rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(rect, BLUE)
    set_no_line(rect)
    rect.adjustments[0] = 0.04

    # Number circle (orange)
    circle_left = left + Emu(146000)
    circle_top = top + Emu(120000)
    add_numbered_circle(slide, circle_left, circle_top, number, fill_color=ORANGE)

    # Title text (bold, white, next to number)
    title_left = circle_left + Emu(340000)
    title_txb = slide.shapes.add_textbox(title_left, circle_top - Emu(20000), width - Emu(530000), Emu(280000))
    tf = title_txb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE

    # Body text (bullets below title)
    body_top = circle_top + Emu(280000)
    body_txb = slide.shapes.add_textbox(circle_left, body_top, width - Emu(290000), height - Emu(420000))
    tf = body_txb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = bullet
        r.font.name = FONT
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)  # light gray on blue bg
        p.space_before = Pt(2)


def add_numbered_item(slide, left, top, number, text, width=Emu(5400000)):
    """Pattern B: Numbered circle + text next to it."""
    add_numbered_circle(slide, left, top, number, fill_color=DARK_NAVY)
    txb = slide.shapes.add_textbox(left + Emu(340000), top - Emu(18000), width, Emu(300000))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(11)
    r.font.color.rgb = DARK_TEXT


def set_cell_bg(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:srgbClr'), {'val': str(color)})
    sf.append(clr)
    tcPr.append(sf)


def style_cell(cell, text, bold=False, color=DARK_TEXT, bg=None, size=Pt(10), align=None):
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(45720)
    tf.margin_right = Emu(45720)
    tf.margin_top = Emu(18288)
    tf.margin_bottom = Emu(18288)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    if bg:
        set_cell_bg(cell, bg)


# ─── LOAD TEMPLATE ─────────────────────────────────────────────────────
print("Loading template...")
prs = Presentation(r'C:\Users\2750834\Downloads\Haryana.pptx')

# Find layout
layout = None
for l in prs.slide_layouts:
    if l.name == '1_Title and Content':
        layout = l
        break
if not layout:
    for l in prs.slide_layouts:
        if l.name == 'Blank':
            layout = l
            break

# Delete existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

print("Template cleared. Building 4 slides...\n")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: FY25-26 Review
# Layout: Two-column with section banners + numbered items
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(layout)
add_title(slide, "FY25-26 Review")
add_slide_number(slide, 1)

# ── LEFT: What Went Well ──
left_x = Emu(250000)
banner_w = Emu(5700000)
add_section_banner(slide, left_x, Emu(850000), banner_w, Emu(340000), "What Went Well", GREEN_ACCENT)

items_well = [
    "Dealer engagement schemes and flat PL supported sales",
    "Q1 & Q2 showed good demand and NOD performance",
    "Delivery time improved from 4-5 days to 1-2 days",
    "Annual dealer tour was motivating for channel partners",
]

y = Emu(1320000)
for i, item in enumerate(items_well):
    add_numbered_item(slide, left_x + Emu(70000), y, i + 1, item, width=Emu(5200000))
    y += Emu(420000)

# ── RIGHT: Areas for Improvement ──
right_x = Emu(6250000)
add_section_banner(slide, right_x, Emu(850000), banner_w, Emu(340000), "Areas for Improvement")

items_issues = [
    "8mm/10mm price gap higher than industry; frequent pricing changes hurt dealer confidence",
    "JSW One TMT brand positioning unclear; intra-brand competition with Neo affected conversions",
    "550D non-availability in H1 and product shining below premium brand standards",
    "Pricing shift from premium-economical to purely premium reduced market competitiveness",
]

y = Emu(1320000)
for i, item in enumerate(items_issues):
    add_numbered_item(slide, right_x + Emu(70000), y, i + 1, item, width=Emu(5200000))
    y += Emu(500000)

# ── BOTTOM: Divider line ──
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(250000), Emu(6100000), Emu(11700000), Emu(3000))
set_fill(line, RGBColor(0xD1, 0xD5, 0xDB))
set_no_line(line)

print("  Slide 1: FY25-26 Review ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: Growth Strategy & Helix TMT Launch
# Layout: LEFT = text + table | RIGHT = 4 blue numbered blocks
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(layout)
add_title(slide, "Growth Strategy & Helix TMT Launch")
add_slide_number(slide, 2)

# ── LEFT side: Subtitle + Channel Criteria Table ──
sub = slide.shapes.add_textbox(Emu(250000), Emu(850000), Emu(7000000), Emu(500000))
tf = sub.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "2X volume growth through SOB expansion + new dealer appointments in white-space markets"
r.font.name = FONT
r.font.size = Pt(11)
r.font.color.rgb = DARK_TEXT

# Section label
add_section_banner(slide, Emu(250000), Emu(1420000), Emu(7000000), Emu(300000), "Helix Channel Selection Criteria")

# Table
table_data = [
    ["Scenario", "Action Plan"],
    ["High-Performing\nJSW Dealers", "Onboard for Helix only if high financial capacity and untapped market potential"],
    ["Limited Capacity\nDealers", "Appoint new/alternative partners to ensure Helix gets dedicated focus"],
    ["Target Counters", "Focus on retailers moving high volumes of Jindal/Rapid/Rungta"],
]
tbl_shape = slide.shapes.add_table(4, 2, Emu(250000), Emu(1820000), Emu(7000000), Emu(1800000))
tbl = tbl_shape.table
tbl.columns[0].width = Emu(2100000)
tbl.columns[1].width = Emu(4900000)

for ri, row in enumerate(table_data):
    for ci, txt in enumerate(row):
        cell = tbl.cell(ri, ci)
        if ri == 0:
            style_cell(cell, txt, bold=True, color=WHITE, bg=BLUE, size=Pt(10))
        else:
            bg = TABLE_ALT if ri % 2 == 0 else None
            style_cell(cell, txt, bold=(ci == 0), color=DARK_TEXT, bg=bg, size=Pt(10))

# Pricing note below table
pnote = slide.shapes.add_textbox(Emu(250000), Emu(3750000), Emu(7000000), Emu(700000))
tf = pnote.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Pricing: "
r.font.name = FONT
r.font.size = Pt(10)
r.font.bold = True
r.font.color.rgb = DARK_TEXT
r = tf.paragraphs[0].add_run()
r.text = "Align with secondary brands (Jindal, Rapid, Rungta). Position as \"Best Product at Good Competitive Price\" without hitting premium JSW levels."
r.font.name = FONT
r.font.size = Pt(10)
r.font.color.rgb = DARK_TEXT

p2 = tf.add_paragraph()
p2.space_before = Pt(6)
r2 = p2.add_run()
r2.text = "Stocking: "
r2.font.name = FONT
r2.font.size = Pt(10)
r2.font.bold = True
r2.font.color.rgb = DARK_TEXT
r2 = p2.add_run()
r2.text = "No yard stock needed currently — direct plant supply is cost-competitive. Re-evaluate if required."
r2.font.name = FONT
r2.font.size = Pt(10)
r2.font.color.rgb = DARK_TEXT

# ── RIGHT side: 4 blue numbered blocks ──
block_x = Emu(7600000)
block_w = Emu(4400000)
block_h = Emu(1100000)
gap = Emu(120000)

blocks = [
    ("Increase SOB of existing channels", [
        "Strengthen volume per dealer",
        "Improve monthly transacting frequency",
    ]),
    ("New dealers in white-space markets", [
        "Target low premium-brand acceptance areas",
        "Focus where secondary brands dominate",
    ]),
    ("Helix: Economy segment positioning", [
        "Rural markets with low primary steel demand",
        "Customers prioritizing quality at competitive price",
    ]),
    ("Brand differentiation", [
        "JSW One → Premium, full services, marketing",
        "One Helix → Value-focused, wider channel",
    ]),
]

y = Emu(850000)
for i, (title, bullets) in enumerate(blocks):
    add_blue_block(slide, block_x, y, block_w, block_h, i + 1, title, bullets)
    y += block_h + gap

print("  Slide 2: Growth Strategy & Helix Launch ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: Industry Practices & Market Activation
# Layout: TOP = best practices (two-column) | BOTTOM = meets plan + table
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(layout)
add_title(slide, "Industry Practices & Market Activation Plan")
add_slide_number(slide, 3)

# ── TOP: Industry Best Practices ──
add_section_banner(slide, Emu(250000), Emu(850000), Emu(11700000), Emu(300000), "Industry Best Practices")

# Left: Primary brands
pleft = slide.shapes.add_textbox(Emu(250000), Emu(1260000), Emu(5700000), Emu(1200000))
tf = pleft.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Primary Brands: TATA, Jindal Panther"
r.font.name = FONT
r.font.size = Pt(11)
r.font.bold = True
r.font.color.rgb = BLUE

items = ["Govt approvals (PWD)", "Influencer loyalty (Architects)", "Strong brand pull"]
for item in items:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = f"• {item}"
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = DARK_TEXT

# Right: Secondary brands
pright = slide.shapes.add_textbox(Emu(6250000), Emu(1260000), Emu(5700000), Emu(1200000))
tf = pright.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Secondary Brands: Jindal, Rapid/Rungta"
r.font.name = FONT
r.font.size = Pt(11)
r.font.bold = True
r.font.color.rgb = BLUE

items = ['"Sauda" booking culture', "Transparent schemes", "Pure price play", "High rural availability / small supplies"]
for item in items:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = f"• {item}"
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = DARK_TEXT

# ── BOTTOM: Influencer & Channel Meets Plan ──
add_section_banner(slide, Emu(250000), Emu(2800000), Emu(11700000), Emu(300000), "Influencer & Channel Meets Plan (FY26-27)")

# Key points
y_start = Emu(3230000)
items_meets = [
    "Nuh & Palwal — Very low branded steel sale; plan One Helix for these markets",
    "NCR Focus — Architects, engineers, builders are decision makers; plan effective activities",
    "Requirement: 2 TSEs for Gautam Budh Nagar + Hapur; 1 DSR for Project sales",
    "Influencer meets on demand only — should not be linked to target or incentive",
]
y = y_start
for i, item in enumerate(items_meets):
    add_numbered_item(slide, Emu(320000), y, i + 1, item, width=Emu(7200000))
    y += Emu(420000)

# Meet frequency table (right side of bottom section)
mtbl_data = [
    ["Meet Type", "Frequency", "Level"],
    ["Mason Meets", "Every 2 months", "At counters"],
    ["Contractor Meets", "Every 2 months", "District level"],
    ["Architect Meetings", "Half-yearly / Yearly", "Regional"],
]
tbl_shape = slide.shapes.add_table(4, 3, Emu(7800000), Emu(3230000), Emu(4100000), Emu(1300000))
tbl = tbl_shape.table
tbl.columns[0].width = Emu(1500000)
tbl.columns[1].width = Emu(1400000)
tbl.columns[2].width = Emu(1200000)

for ri, row in enumerate(mtbl_data):
    for ci, txt in enumerate(row):
        cell = tbl.cell(ri, ci)
        if ri == 0:
            style_cell(cell, txt, bold=True, color=WHITE, bg=BLUE, size=Pt(9))
        else:
            bg = TABLE_ALT if ri % 2 == 0 else None
            style_cell(cell, txt, color=DARK_TEXT, bg=bg, size=Pt(9))

# North Haryana note
note = slide.shapes.add_textbox(Emu(320000), Emu(4950000), Emu(11600000), Emu(500000))
tf = note.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "North Haryana: "
r.font.name = FONT
r.font.size = Pt(10)
r.font.bold = True
r.font.color.rgb = DARK_TEXT
r = tf.paragraphs[0].add_run()
r.text = "Mason meets at counters + contractor meets at district level every 2 months; architect meetings half-yearly/yearly — strengthens product awareness, generates leads, improves project conversions."
r.font.name = FONT
r.font.size = Pt(10)
r.font.color.rgb = DARK_TEXT

print("  Slide 3: Industry Practices & Market Activation ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: Distributor Performance KPIs (FY26-27)
# Layout: 3 large blue blocks, full width, stacked
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(layout)
add_title(slide, "Distributor Key Performance Improvements (FY26-27)")
add_slide_number(slide, 4)

# Subtitle
sub = slide.shapes.add_textbox(Emu(250000), Emu(780000), Emu(11700000), Emu(300000))
tf = sub.text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Three key performance points distributors must improve to meet business expectations"
r.font.name = FONT
r.font.size = Pt(11)
r.font.color.rgb = DARK_TEXT

block_x = Emu(250000)
block_w = Emu(11700000)
block_h = Emu(1450000)
gap = Emu(150000)
y = Emu(1200000)

kpi_blocks = [
    ("Expand Market Reach & Sales Volume", [
        "Increase retailer coverage and actively engage with contractors and construction projects",
        "Drive higher demand through wider channel penetration across districts",
    ]),
    ("Appoint Dedicated Branch Manager", [
        "Ensure focused leadership at branch level to drive sales execution and team management",
        "Enable faster decision-making and accountability at the local level",
    ]),
    ("Improve Joint Visits with JSW One Team", [
        "Strengthen collaboration with channel partners, key influencers, and architects",
        "Regular joint visits to improve brand visibility and deepen relationships",
    ]),
]

for i, (title, bullets) in enumerate(kpi_blocks):
    add_blue_block(slide, block_x, y, block_w, block_h, i + 1, title, bullets)
    y += block_h + gap

print("  Slide 4: Distributor KPIs ✓")

# ─── SAVE ───────────────────────────────────────────────────────────────
output = r'D:\RandomTestsClaude\Haryana_NewSlides.pptx'
prs.save(output)
print(f"\nSaved: {output}")
print(f"Total: {len(prs.slides)} slides")
