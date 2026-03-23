import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── COLORS ─────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x18, 0x48, 0x9D)
DARK_NAVY  = RGBColor(0x0D, 0x2B, 0x5E)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
GRAY_TITLE  = RGBColor(0x7F, 0x7F, 0x7F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE      = RGBColor(0xD9, 0x77, 0x06)
LIGHT_SLATE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN       = RGBColor(0x05, 0x72, 0x3A)
TABLE_ALT   = RGBColor(0xF0, 0xF4, 0xF8)
LIGHT_BLUE  = RGBColor(0x93, 0xC5, 0xFD)

FONT = 'Calibri'

# ─── Margins (in EMU) ───────────────────────────────────────────────────
MARGIN_L    = Emu(330000)   # ~0.36in left margin
FULL_W      = Emu(11530000) # ~12.61in usable width
HALF_W      = Emu(5550000)  # ~6.07in per column
COL_GAP     = Emu(430000)   # ~0.47in gap between columns
RIGHT_COL_X = MARGIN_L + HALF_W + COL_GAP  # ~6.90in


# ─── HELPER FUNCTIONS ────────────────────────────────────────────────────

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_title(slide, text):
    """Gray 20pt bold title at top of slide."""
    txb = slide.shapes.add_textbox(Emu(160000), Emu(200000), Emu(12060000), Emu(380000))
    tf = txb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = GRAY_TITLE

def add_slide_num(slide, n):
    txb = slide.shapes.add_textbox(Emu(9448800), Emu(6484939), Emu(2743200), Emu(365125))
    tf = txb.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    r = tf.paragraphs[0].add_run()
    r.text = str(n)
    r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = GRAY_TITLE

def add_navy_bar(slide, text, y=Emu(6050000)):
    """Dark navy footer bar with white text (like main file's KPI bar)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, Emu(12192000), Emu(420000))
    set_fill(bar, DARK_NAVY); no_line(bar)
    tf = bar.text_frame
    tf.margin_left = Emu(400000); tf.margin_top = Emu(50000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = WHITE

def add_section_banner(slide, x, y, w, text, fill=BLUE, h=Emu(310000)):
    """Rounded rectangle section header with white text."""
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(rect, fill); no_line(rect)
    rect.adjustments[0] = 0.06
    tf = rect.text_frame
    tf.margin_left = Emu(100000); tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE

def add_numbered_circle(slide, x, y, num, fill=DARK_NAVY, sz=Emu(250000)):
    """Small numbered circle."""
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    set_fill(o, fill); no_line(o)
    tf = o.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(num)
    r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = WHITE

def add_numbered_item(slide, x, y, num, text, w=Emu(5200000)):
    """Circle + adjacent text on same line."""
    add_numbered_circle(slide, x, y, num)
    txb = slide.shapes.add_textbox(x + Emu(320000), y - Emu(10000), w, Emu(280000))
    tf = txb.text_frame; tf.word_wrap = True; tf.margin_left = 0; tf.margin_top = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = DARK_TEXT

def add_blue_block(slide, x, y, w, h, num, title, bullets):
    """
    Blue rounded rect with ALL text inside its own text frame.
    Only the orange number circle is a separate overlapping shape.
    """
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(rect, BLUE); no_line(rect)
    rect.adjustments[0] = 0.04
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(430000)  # leave room for the number circle
    tf.margin_right = Emu(80000)
    tf.margin_top = Emu(60000)
    tf.margin_bottom = Emu(40000)

    # Title paragraph
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE

    # Bullet paragraphs
    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = f"• {b}"
        r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = LIGHT_SLATE

    # Orange number circle (only overlapping element — intentional)
    add_numbered_circle(slide, x + Emu(80000), y + Emu(55000), num, fill=ORANGE)

def set_cell_bg(cell, color):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    sf = tcPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': str(color)}))
    tcPr.append(sf)

def style_cell(cell, text, bold=False, color=DARK_TEXT, bg=None, size=Pt(10)):
    cell.text = ''
    tf = cell.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(45720); tf.margin_right = Emu(45720)
    tf.margin_top = Emu(18288); tf.margin_bottom = Emu(18288)
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.name = FONT; r.font.size = size; r.font.bold = bold; r.font.color.rgb = color
    if bg: set_cell_bg(cell, bg)

def add_bullet_textbox(slide, x, y, w, h, header, items, header_color=BLUE):
    """Text box with a bold header and bullet items below."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame; tf.word_wrap = True; tf.margin_left = 0; tf.margin_top = 0

    # Header
    r = tf.paragraphs[0].add_run()
    r.text = header
    r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = header_color

    # Bullets
    for item in items:
        p = tf.add_paragraph(); p.space_before = Pt(4)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = DARK_TEXT
    return txb


# ─── LOAD TEMPLATE ──────────────────────────────────────────────────────
print("Loading template...")
prs = Presentation(r'C:\Users\2750834\Downloads\Haryana.pptx')

# Use Blank layout — no phantom placeholders
blank_layout = None
for l in prs.slide_layouts:
    if l.name == 'Blank':
        blank_layout = l; break

# Delete existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

print("Building 4 slides with Blank layout...\n")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: FY25-26 Review
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_title(s, "FY25-26 Review")
add_slide_num(s, 1)

# Left column: What Went Well
add_section_banner(s, MARGIN_L, Emu(780000), HALF_W, "What Went Well", fill=GREEN)

well_items = [
    "Dealer engagement schemes and flat PL supported sales",
    "Q1 & Q2 showed good demand and NOD performance",
    "Delivery time improved from 4\u20135 days to 1\u20132 days",
    "Annual dealer tour was motivating for channel partners",
]
y = Emu(1220000)
for i, item in enumerate(well_items):
    add_numbered_item(s, MARGIN_L + Emu(60000), y, i+1, item, w=HALF_W - Emu(400000))
    y += Emu(550000)

# Right column: Areas for Improvement
add_section_banner(s, RIGHT_COL_X, Emu(780000), HALF_W, "Areas for Improvement")

issue_items = [
    "8mm/10mm price gap above industry standard; frequent pricing changes hurt dealer confidence",
    "JSW One TMT brand positioning unclear; intra-brand competition with Neo affected conversions",
    "550D non-availability in H1; product shining below premium brand expectations",
    "Pricing shift from premium-economical to purely premium reduced competitiveness",
]
y = Emu(1220000)
for i, item in enumerate(issue_items):
    add_numbered_item(s, RIGHT_COL_X + Emu(60000), y, i+1, item, w=HALF_W - Emu(400000))
    y += Emu(600000)

# Footer bar
add_navy_bar(s, "Focus: Stabilize pricing strategy, strengthen brand positioning, and improve product shining to match premium standards")
print("  Slide 1: FY25-26 Review ✓")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: Growth Strategy & Helix TMT Launch
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_title(s, "Growth Strategy & Helix TMT Launch")
add_slide_num(s, 2)

# LEFT: Summary + table + notes
left_w = Emu(7100000)
sub = slide_shapes = s.shapes.add_textbox(MARGIN_L, Emu(780000), left_w, Emu(350000))
tf = sub.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "2X volume growth through SOB expansion + new dealer appointments in white-space markets"
r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = DARK_TEXT

add_section_banner(s, MARGIN_L, Emu(1230000), left_w, "Helix Channel Selection Criteria")

# Table
tbl_data = [
    ["Scenario", "Action Plan"],
    ["High-Performing JSW Dealers", "Onboard for Helix only if high financial capacity and untapped market potential"],
    ["Limited Capacity Dealers", "Appoint new/alternative partners to ensure Helix gets dedicated focus"],
    ["Target Counters", "Focus on retailers moving high volumes of Jindal/Rapid/Rungta"],
]
ts = s.shapes.add_table(4, 2, MARGIN_L, Emu(1640000), left_w, Emu(1500000))
t = ts.table
t.columns[0].width = Emu(2200000); t.columns[1].width = Emu(4900000)
for ri, row in enumerate(tbl_data):
    for ci, txt in enumerate(row):
        c = t.cell(ri, ci)
        if ri == 0:
            style_cell(c, txt, bold=True, color=WHITE, bg=BLUE, size=Pt(10))
        else:
            style_cell(c, txt, bold=(ci==0), color=DARK_TEXT, bg=TABLE_ALT if ri%2==0 else None, size=Pt(10))

# Notes below table
notes = s.shapes.add_textbox(MARGIN_L, Emu(3260000), left_w, Emu(800000))
tf = notes.text_frame; tf.word_wrap = True

p = tf.paragraphs[0]
r = p.add_run(); r.text = "Pricing: "; r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = DARK_TEXT
r = p.add_run(); r.text = 'Align with Jindal, Rapid, Rungta. "Best Product at Good Competitive Price" without hitting premium JSW levels.'; r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = DARK_TEXT

p = tf.add_paragraph(); p.space_before = Pt(6)
r = p.add_run(); r.text = "Stocking: "; r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = DARK_TEXT
r = p.add_run(); r.text = "No yard stock needed \u2014 direct plant supply is cost-competitive. Re-evaluate if required."; r.font.name = FONT; r.font.size = Pt(10); r.font.color.rgb = DARK_TEXT

# RIGHT: 4 blue blocks
bx = MARGIN_L + left_w + Emu(200000)  # ~7.63in
bw = Emu(4570000)  # ~5.0in
bh = Emu(1050000)
bg = Emu(100000)
blocks = [
    ("Increase SOB of existing channels", ["Strengthen volume per dealer", "Improve monthly transacting frequency"]),
    ("New dealers in white-space markets", ["Target low premium-brand acceptance areas", "Focus where secondary brands dominate"]),
    ("Helix: Economy segment positioning", ["Rural markets with low primary steel demand", "Quality at competitive price, not premium"]),
    ("Brand differentiation", ["JSW One \u2192 Premium, full services & marketing", "One Helix \u2192 Value-focused, wider channel"]),
]
y = Emu(780000)
for i, (title, bullets) in enumerate(blocks):
    add_blue_block(s, bx, y, bw, bh, i+1, title, bullets)
    y += bh + bg

# Footer
add_navy_bar(s, "Target: Economy segment in rural markets through new channel appointments and competitive pricing")
print("  Slide 2: Growth Strategy & Helix Launch ✓")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: Industry Practices & Market Activation
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_title(s, "Industry Practices & Market Activation Plan")
add_slide_num(s, 3)

# TOP: Best Practices banner + two columns
add_section_banner(s, MARGIN_L, Emu(780000), FULL_W, "Industry Best Practices")

add_bullet_textbox(s, MARGIN_L, Emu(1200000), HALF_W, Emu(1100000),
    "Primary Brands: TATA, Jindal Panther",
    ["Govt approvals (PWD)", "Influencer loyalty (Architects)", "Strong brand pull"])

add_bullet_textbox(s, RIGHT_COL_X, Emu(1200000), HALF_W, Emu(1100000),
    "Secondary Brands: Jindal, Rapid/Rungta",
    ['"Sauda" booking culture', "Transparent schemes", "Pure price play", "High rural availability / small supplies"])

# BOTTOM: Meets Plan banner + items + table
add_section_banner(s, MARGIN_L, Emu(2600000), FULL_W, "Influencer & Channel Meets Plan (FY26-27)")

meet_items = [
    "Nuh & Palwal \u2014 Very low branded steel sale; plan One Helix",
    "NCR \u2014 Architects, engineers, builders are decision makers",
    "Need: 2 TSEs (Gautam Budh Nagar + Hapur) + 1 DSR (Projects)",
    "Influencer meets: On demand only, not linked to targets",
]
y = Emu(3050000)
for i, item in enumerate(meet_items):
    add_numbered_item(s, MARGIN_L + Emu(60000), y, i+1, item, w=Emu(6700000))
    y += Emu(430000)

# Frequency table (right side)
mt_data = [
    ["Meet Type", "Frequency", "Level"],
    ["Mason Meets", "Every 2 months", "Counters"],
    ["Contractor Meets", "Every 2 months", "District"],
    ["Architect Meets", "Half-yearly", "Regional"],
]
ts = s.shapes.add_table(4, 3, Emu(7900000), Emu(3050000), Emu(3960000), Emu(1200000))
t = ts.table
t.columns[0].width = Emu(1400000); t.columns[1].width = Emu(1360000); t.columns[2].width = Emu(1200000)
for ri, row in enumerate(mt_data):
    for ci, txt in enumerate(row):
        c = t.cell(ri, ci)
        if ri == 0: style_cell(c, txt, bold=True, color=WHITE, bg=BLUE, size=Pt(9))
        else: style_cell(c, txt, color=DARK_TEXT, bg=TABLE_ALT if ri%2==0 else None, size=Pt(9))

# Summary note
note = s.shapes.add_textbox(MARGIN_L, Emu(4800000), FULL_W, Emu(400000))
tf = note.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "North Haryana: Regular meets at counters and district level strengthen awareness, generate leads, and improve conversions."
r.font.name = FONT; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = DARK_TEXT
print("  Slide 3: Industry Practices & Market Activation ✓")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: Distributor KPIs
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_title(s, "Distributor Key Performance Improvements (FY26-27)")
add_slide_num(s, 4)

sub = s.shapes.add_textbox(MARGIN_L, Emu(730000), FULL_W, Emu(280000))
tf = sub.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Three key performance points distributors must improve to meet business expectations"
r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = DARK_TEXT

bx = MARGIN_L; bw = FULL_W; bh = Emu(1150000); bg = Emu(120000)
kpis = [
    ("Expand Market Reach & Sales Volume", [
        "Increase retailer coverage across districts",
        "Actively engage with contractors and construction projects for higher demand",
    ]),
    ("Appoint Dedicated Branch Manager", [
        "Ensure focused leadership at branch level for sales execution",
        "Enable faster decision-making and accountability locally",
    ]),
    ("Improve Joint Visits with JSW One Team", [
        "Strengthen collaboration with channel partners, key influencers, and architects",
        "Regular joint visits to improve brand visibility and relationships",
    ]),
]
y = Emu(1150000)
for i, (title, bullets) in enumerate(kpis):
    add_blue_block(s, bx, y, bw, bh, i+1, title, bullets)
    y += bh + bg

add_navy_bar(s, "Priority: Wider retailer reach, dedicated branch leadership, and closer collaboration with JSW One team")
print("  Slide 4: Distributor KPIs ✓")


# ─── SAVE ────────────────────────────────────────────────────────────────
out = r'D:\RandomTestsClaude\Haryana_NewSlides.pptx'
prs.save(out)
print(f"\nSaved: {out}")
print(f"Total: {len(prs.slides)} slides")
