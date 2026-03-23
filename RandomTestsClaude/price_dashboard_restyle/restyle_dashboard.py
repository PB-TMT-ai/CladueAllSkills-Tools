import sys, copy
sys.stdout.reconfigure(encoding='utf-8')

import openpyxl
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# ─── PATHS ──────────────────────────────────────────────────────────────
INPUT_PPT  = r'D:\JSW_TMT_Price_Dashboard.pptx'
EXCEL_PATH = r'D:\RandomTestsClaude\North TMT pricing_9th Mar.xlsx'
OUTPUT_PPT = r'D:\RandomTestsClaude\price_dashboard_restyle\JSW_TMT_Price_Dashboard_Restyled.pptx'
SHEET_NAME = '1_Summary at PL'

# ─── FILL COLOR MAPPING (old hex → new hex) ─────────────────────────────
# NOTE: E8272A is handled contextually (branding vs status) — see map_shape_fill
FILL_MAP = {
    '1A2B3C': 'F0F4F8',   # dark card bg → light gray
    '0D1B2A': 'FFFFFF',   # darker bg → white
    '223447': '18489D',   # section header → blue
    '2A2A18': 'F0F4F8',   # TISCON row → light banding
    '2A1518': 'F0F4F8',   # JSW row → light banding
    '1A6FB5': '18489D',   # summary accent → blue
    'F5A623': 'D97706',   # amber → Haryana amber
    '27AE60': '05723A',   # green → Haryana green
}

# Red #E8272A context: branding → blue, status → red
RED_BRAND  = '18489D'   # JSW red branding → Haryana blue
RED_STATUS = 'DC2626'   # Action Needed status → stays red

# Fills where text should remain WHITE (dark backgrounds)
DARK_FILLS = {'18489D', '0D2B5E', '05723A', 'D97706', 'DC2626'}

# Text color mapping (non-white, non-red colors)
TEXT_COLOR_MAP = {
    '8BA3C0': '7F7F7F',   # muted blue → gray
    'F5A623': 'D97706',   # amber → Haryana amber
    '27AE60': '05723A',   # green → Haryana green
    '1A6FB5': '18489D',   # blue accent → Haryana blue
}
# E8272A text is handled contextually in restyle_text_runs

TARGET_FONT = 'Calibri'

# ─── SLIDE ↔ MARKET MAPPING ─────────────────────────────────────────────
SLIDE_MARKET_MAP = {
    2: 'DL', 3: 'HR', 4: 'RJ', 5: 'UP', 6: 'CH',
    7: 'PB', 8: 'UK', 9: 'HP', 10: 'JK',
}

MARKET_COLUMNS = {
    'DL': {'distr_col': 2,  'dealer_col': 3,  'gap_col': 4},
    'HR': {'distr_col': 5,  'dealer_col': 6,  'gap_col': 7},
    'RJ': {'distr_col': 8,  'dealer_col': 9,  'gap_col': 10},
    'UP': {'distr_col': 11, 'dealer_col': 12, 'gap_col': 13},
    'CH': {'distr_col': 14, 'dealer_col': 15, 'gap_col': 16},
    'PB': {'distr_col': 17, 'dealer_col': 18, 'gap_col': 19},
    'UK': {'distr_col': 20, 'dealer_col': 21, 'gap_col': 22},
    'HP': {'distr_col': 23, 'dealer_col': 24, 'gap_col': 25},
    'JK': {'distr_col': 26, 'dealer_col': 27, 'gap_col': 28},
}

BRAND_ROWS = range(4, 13)
SKIP_ROWS  = {9}


# ─── HELPER FUNCTIONS ───────────────────────────────────────────────────

def hex_to_rgb(h):
    """Convert hex string like 'E8272A' to RGBColor."""
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_fill_hex(shape):
    """Get the solid fill color hex of a shape, or None."""
    try:
        if shape.fill.type is not None:
            return str(shape.fill.fore_color.rgb).upper()
    except:
        pass
    return None


def _is_status_shape(shape):
    """Check if a red shape is used for status (ACTION NEEDED) vs branding."""
    if shape.has_text_frame:
        text = shape.text_frame.text.upper()
        if 'ACTION' in text or 'NEEDED' in text:
            return True
    return False


def map_shape_fill(shape):
    """Map shape's fill color to Haryana palette. Returns new fill hex or None."""
    old_hex = get_fill_hex(shape)
    if not old_hex:
        return old_hex

    # Handle E8272A contextually: branding → blue, status → red
    if old_hex == 'E8272A':
        if _is_status_shape(shape):
            new_hex = RED_STATUS  # keep red for status badges
        else:
            new_hex = RED_BRAND   # blue for branding
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(new_hex)
        return new_hex

    if old_hex in FILL_MAP:
        new_hex = FILL_MAP[old_hex]
        if new_hex == 'FFFFFF':
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(new_hex)
        return new_hex
    return old_hex


def restyle_text_runs(shape, parent_fill_hex):
    """Restyle all text runs in a shape: map fonts and colors."""
    if not shape.has_text_frame:
        return

    # Is the shape background dark? (text should stay white on dark fills)
    is_dark_bg = parent_fill_hex is not None and parent_fill_hex.upper() in DARK_FILLS

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            # ── Font mapping ──
            if run.font.name in ('Arial', 'Trebuchet MS', 'Consolas', None):
                run.font.name = TARGET_FONT

            # ── Color mapping ──
            try:
                old_color = str(run.font.color.rgb).upper() if run.font.color and run.font.color.rgb else None
            except:
                old_color = None

            if old_color is None:
                # No explicit color — set based on background
                if is_dark_bg:
                    run.font.color.rgb = hex_to_rgb('FFFFFF')
                else:
                    run.font.color.rgb = hex_to_rgb('1E293B')
            elif old_color == 'FFFFFF':
                # White text — keep white on dark bg, make dark on light bg
                if not is_dark_bg:
                    run.font.color.rgb = hex_to_rgb('1E293B')
            elif old_color == 'E8272A':
                # Red text: status indicators (gap values with ▼, "ACTION") → red
                # Branding (JSW names, headers) → blue
                txt = run.text.upper()
                if '\u25bc' in run.text or 'ACTION' in txt or 'NEEDED' in txt:
                    run.font.color.rgb = hex_to_rgb(RED_STATUS)
                else:
                    run.font.color.rgb = hex_to_rgb(RED_BRAND)
            elif old_color in TEXT_COLOR_MAP:
                run.font.color.rgb = hex_to_rgb(TEXT_COLOR_MAP[old_color])
            # else: keep the color as-is


def restyle_shape(shape):
    """Restyle a single shape: fill + text."""
    new_fill = map_shape_fill(shape)
    restyle_text_runs(shape, new_fill)


# ─── EXCEL DATA READER ──────────────────────────────────────────────────

def read_market_data(ws, market_key):
    """Read brand pricing data for a single market."""
    cfg = MARKET_COLUMNS[market_key]
    brands = []
    for row in BRAND_ROWS:
        if row in SKIP_ROWS:
            continue
        brand_name = ws.cell(row=row, column=1).value
        if not brand_name or str(brand_name).startswith('<'):
            continue
        distr = ws.cell(row=row, column=cfg['distr_col']).value
        dealer = ws.cell(row=row, column=cfg['dealer_col']).value
        gap   = ws.cell(row=row, column=cfg['gap_col']).value
        if isinstance(distr, str): distr = None
        if isinstance(dealer, str): dealer = None
        if isinstance(gap, str): gap = None
        brands.append({
            'name':  str(brand_name).strip(),
            'distr': distr,
            'dealer': dealer,
            'gap':   gap if gap is not None else 0,
        })
    return brands


# ─── CHART CREATION (native PPT chart, Haryana colors) ──────────────────

def add_native_chart(slide, data, position):
    """Add a grouped bar chart with TISCON reference line at the given position."""
    left, top, width, height = position

    chart_data = CategoryChartData()
    brand_names = [b['name'] for b in data]
    chart_data.categories = brand_names

    distr_values = tuple(b['distr'] or 0 for b in data)
    dealer_values = tuple(b['dealer'] or 0 for b in data)
    tiscon_dealer = data[0]['dealer'] if data[0]['dealer'] else 0
    tiscon_ref = tuple([tiscon_dealer] * len(data))

    chart_data.add_series('Distributor', distr_values)
    chart_data.add_series('Dealer', dealer_values)
    chart_data.add_series('TISCON Ref', tiscon_ref)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = chart_frame.chart

    # Legend
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = TARGET_FONT
    chart.legend.font.size = Pt(7)

    plot = chart.plots[0]
    plot.gap_width = 80

    # Series 0 — Distributor: Blue #18489D
    s0 = plot.series[0]
    s0.format.fill.solid()
    s0.format.fill.fore_color.rgb = hex_to_rgb('18489D')

    # Series 1 — Dealer: Amber #D97706
    s1 = plot.series[1]
    s1.format.fill.solid()
    s1.format.fill.fore_color.rgb = hex_to_rgb('D97706')

    # Value axis
    value_axis = chart.value_axis
    min_val = min(b['distr'] or 99999 for b in data)
    min_val = min(min_val, min(b['dealer'] or 99999 for b in data))
    value_axis.minimum_scale = max(0, int(min_val - 2000))
    value_axis.has_title = False
    value_axis.tick_labels.font.name = TARGET_FONT
    value_axis.tick_labels.font.size = Pt(7)
    value_axis.tick_labels.number_format = '#,##0'

    # Category axis
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.name = TARGET_FONT
    cat_axis.tick_labels.font.size = Pt(6)

    # Convert 3rd series to line overlay via XML
    _convert_series_to_line(chart)


def _convert_series_to_line(chart):
    """Move the 3rd bar series to a line chart overlay via XML manipulation."""
    chart_space = chart._chartSpace
    plot_area = chart_space.find(qn('c:chart')).find(qn('c:plotArea'))
    bar_chart = plot_area.find(qn('c:barChart'))

    if bar_chart is None:
        return

    series_list = bar_chart.findall(qn('c:ser'))
    if len(series_list) < 3:
        return

    target_ser = series_list[2]
    bar_chart.remove(target_ser)

    # Create lineChart element
    line_chart = bar_chart.makeelement(qn('c:lineChart'), {})
    grouping = line_chart.makeelement(qn('c:grouping'), {'val': 'standard'})
    line_chart.append(grouping)
    line_chart.append(target_ser)

    # No markers
    marker = target_ser.makeelement(qn('c:marker'), {})
    symbol = marker.makeelement(qn('c:symbol'), {'val': 'none'})
    marker.append(symbol)
    target_ser.append(marker)

    # Style: red dashed line
    spPr = target_ser.find(qn('c:spPr'))
    if spPr is None:
        spPr = target_ser.makeelement(qn('c:spPr'), {})
        target_ser.append(spPr)

    ln = spPr.makeelement(qn('a:ln'), {'w': '19050'})
    sf = ln.makeelement(qn('a:solidFill'), {})
    clr = sf.makeelement(qn('a:srgbClr'), {'val': 'E8272A'})
    sf.append(clr)
    ln.append(sf)

    prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(prstDash)
    spPr.append(ln)

    # No bar fill (it's a line now)
    noFill = spPr.makeelement(qn('a:noFill'), {})
    spPr.insert(0, noFill)

    # Insert lineChart after barChart
    bar_chart.addnext(line_chart)

    # Copy axis references
    ax_ids = bar_chart.findall(qn('c:axId'))
    for ax in ax_ids:
        line_chart.append(copy.deepcopy(ax))


# ─── MAIN ───────────────────────────────────────────────────────────────

def main():
    print("Loading files...")
    prs = Presentation(INPUT_PPT)
    wb  = openpyxl.load_workbook(EXCEL_PATH, data_only=True)
    ws  = wb[SHEET_NAME]

    total_slides = len(prs.slides)
    print(f"  {total_slides} slides, dimensions: {prs.slide_width} x {prs.slide_height}")

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1

        # ── Set slide background to white ──
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # ── Find chart image to replace (slides 2-10 only) ──
        chart_image = None
        market_key = SLIDE_MARKET_MAP.get(slide_num)

        shapes_to_process = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and market_key:
                chart_image = shape
                continue
            shapes_to_process.append(shape)

        # ── Restyle all non-chart shapes ──
        for shape in shapes_to_process:
            restyle_shape(shape)

        # ── Replace chart image with native chart ──
        if chart_image and market_key:
            pos = (chart_image.left, chart_image.top, chart_image.width, chart_image.height)
            # Remove old picture
            sp_elem = chart_image._element
            sp_elem.getparent().remove(sp_elem)
            # Add native chart
            data = read_market_data(ws, market_key)
            add_native_chart(slide, data, pos)
            print(f"  Slide {slide_num}: restyled + chart replaced ({market_key})")
        else:
            label = "title" if slide_num == 1 else ("summary" if slide_num == 11 else "")
            print(f"  Slide {slide_num}: restyled ({label})")

    prs.save(OUTPUT_PPT)
    print(f"\nSaved: {OUTPUT_PPT}")
    print(f"Total: {total_slides} slides")


if __name__ == '__main__':
    main()
