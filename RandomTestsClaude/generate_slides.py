import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

# ── Design constants (matched from main file) ──────────────────────────
TITLE_LEFT = Emu(160412)
TITLE_TOP = Emu(260551)
TITLE_WIDTH = Emu(12057374)
TITLE_HEIGHT = Emu(369332)
TITLE_FONT = 'Calibri'
TITLE_SIZE = Pt(20)
TITLE_COLOR = RGBColor(0x7F, 0x7F, 0x7F)

BODY_FONT = 'Calibri'
BODY_SIZE = Pt(12)
BODY_COLOR = RGBColor(0x1E, 0x29, 0x3B)  # dark slate from main file

SECTION_COLOR = RGBColor(0x18, 0x48, 0x9D)  # blue headers
ACCENT_COLOR = RGBColor(0x18, 0x47, 0xA1)  # title page blue

CONTENT_LEFT = Emu(332220)
CONTENT_TOP = Emu(900000)
CONTENT_WIDTH = Emu(11524664)

SLIDE_NUM_LEFT = Emu(9448800)
SLIDE_NUM_TOP = Emu(6484939)
SLIDE_NUM_WIDTH = Emu(2743200)
SLIDE_NUM_HEIGHT = Emu(365125)

# Table colors
TABLE_HEADER_BG = RGBColor(0x18, 0x48, 0x9D)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xF8)
TABLE_BORDER = RGBColor(0xD1, 0xD5, 0xDB)


def add_title(slide, text):
    """Add a title textbox matching main file design."""
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = TITLE_FONT
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    return txBox


def add_slide_number(slide, num):
    """Add slide number in bottom-right corner."""
    txBox = slide.shapes.add_textbox(SLIDE_NUM_LEFT, SLIDE_NUM_TOP, SLIDE_NUM_WIDTH, SLIDE_NUM_HEIGHT)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.name = TITLE_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)


def add_body_text(slide, left, top, width, height, paragraphs):
    """
    Add a text box with multiple paragraphs.
    Each paragraph is a dict: {text, bold, color, size, bullet, indent}
    or a string for simple text.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, para_data in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(para_data, str):
            para_data = {'text': para_data}

        # Handle mixed runs (list of dicts for mixed formatting)
        if 'runs' in para_data:
            for run_data in para_data['runs']:
                run = p.add_run()
                run.text = run_data.get('text', '')
                run.font.name = BODY_FONT
                run.font.size = run_data.get('size', BODY_SIZE)
                run.font.bold = run_data.get('bold', False)
                run.font.color.rgb = run_data.get('color', BODY_COLOR)
        else:
            run = p.add_run()
            run.text = para_data.get('text', '')
            run.font.name = BODY_FONT
            run.font.size = para_data.get('size', BODY_SIZE)
            run.font.bold = para_data.get('bold', False)
            run.font.color.rgb = para_data.get('color', BODY_COLOR)

        if para_data.get('spacing_before'):
            p.space_before = para_data['spacing_before']
        if para_data.get('spacing_after'):
            p.space_after = para_data['spacing_after']
        if para_data.get('level'):
            p.level = para_data['level']
        if para_data.get('alignment'):
            p.alignment = para_data['alignment']

    return txBox


def add_section_header(paragraphs_list, text, spacing_before=Pt(12)):
    """Helper: add a blue bold section header to a paragraph list."""
    paragraphs_list.append({
        'text': text,
        'bold': True,
        'color': SECTION_COLOR,
        'size': Pt(13),
        'spacing_before': spacing_before
    })


def add_bullet(paragraphs_list, text, bold=False):
    """Helper: add a bullet point to a paragraph list."""
    paragraphs_list.append({
        'text': f'• {text}',
        'bold': bold,
        'size': BODY_SIZE,
        'spacing_before': Pt(4)
    })


def set_cell_bg(cell, color):
    """Set background color of a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': str(color)})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def style_cell(cell, text, font_name='Calibri', font_size=Pt(11), bold=False, color=None, bg=None, align=None):
    """Style a table cell with text and formatting."""
    cell.text = ''
    tf = cell.text_frame
    tf.word_wrap = True
    # Remove default margins for tighter fit
    tf.margin_left = Emu(45720)
    tf.margin_right = Emu(45720)
    tf.margin_top = Emu(27432)
    tf.margin_bottom = Emu(27432)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if bg:
        set_cell_bg(cell, bg)


# ── Load template and prepare ──────────────────────────────────────────
print("Loading main file as template...")
prs = Presentation(r'C:\Users\2750834\Downloads\Haryana.pptx')

# Find the 1_Title and Content layout
target_layout = None
for layout in prs.slide_layouts:
    if layout.name == '1_Title and Content':
        target_layout = layout
        break
if not target_layout:
    # Fallback to Blank
    for layout in prs.slide_layouts:
        if layout.name == 'Blank':
            target_layout = layout
            break

print(f"Using layout: {target_layout.name}")

# Delete all existing slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn('r:id'))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

print("Cleared existing slides. Creating new slides...")

slide_num_counter = 1

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: FY25-26 Review (merged: What Went Well + What Went Not Well)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "FY25-26 Review: What Went Well & Areas for Improvement")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

# Left column - What went well
left_paras = []
add_section_header(left_paras, "What Went Well", spacing_before=Pt(0))
add_bullet(left_paras, "Dealer engagement schemes and flat PL supported sales")
add_bullet(left_paras, "Q1 and part of Q2 showed good demand and NOD performance")
add_bullet(left_paras, "Product shining improved; delivery/service time reduced from 4-5 days to 1-2 days")
add_bullet(left_paras, "Annual dealer tour was motivating for channel partners")

add_body_text(slide,
    Emu(332220), Emu(900000),
    Emu(5500000), Emu(5200000),
    left_paras)

# Right column - What went not well
right_paras = []
add_section_header(right_paras, "Areas for Improvement", spacing_before=Pt(0))
add_bullet(right_paras, "8mm & 10mm dia price difference was higher than industry standards")
add_bullet(right_paras, "Frequent pricing changes impacted dealer confidence and caused losses")
add_bullet(right_paras, "JSW One TMT brand positioning was not clearly established in the market")
add_bullet(right_paras, "Intra-brand competition with Neo affected dealers and site conversions")
add_bullet(right_paras, "Initial non-availability and later uncompetitive pricing of 550D impacted sales (H1 2025)")
add_bullet(right_paras, "Product shining/quality did not fully meet dealer expectations vs premium brands")
add_bullet(right_paras, "Pricing inconsistency reduced competitiveness as brand shifted to purely premium")

add_body_text(slide,
    Emu(6200000), Emu(900000),
    Emu(5600000), Emu(5200000),
    right_paras)

print("  Slide 1: FY25-26 Review ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: Dealer Growth Plan for 2X Volume
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Dealer Growth Plan for 2X Volume")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

paras = []
add_section_header(paras, "District-wise Active Dealer Planning, New Appointments & Timelines", spacing_before=Pt(0))
paras.append({'text': '', 'size': Pt(6)})  # spacer

paras.append({
    'runs': [
        {'text': 'Volume will be achieved through ', 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'two approaches:', 'size': BODY_SIZE, 'bold': True, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})

paras.append({'text': '', 'size': Pt(6)})

# Approach 1
paras.append({
    'runs': [
        {'text': '1', 'size': Pt(14), 'bold': True, 'color': SECTION_COLOR},
        {'text': '  Increasing the SOB of existing channels', 'size': Pt(13), 'bold': True, 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(12)
})
add_bullet(paras, "Strengthen Share of Business with current dealer network")
add_bullet(paras, "Improve monthly transacting frequency and volume per dealer")

paras.append({'text': '', 'size': Pt(6)})

# Approach 2
paras.append({
    'runs': [
        {'text': '2', 'size': Pt(14), 'bold': True, 'color': SECTION_COLOR},
        {'text': '  Appointing new dealers in white-space markets', 'size': Pt(13), 'bold': True, 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(12)
})
add_bullet(paras, "Target markets where premium brand acceptance is low")
add_bullet(paras, "Focus on areas where secondary brands (Jindal, Rapid, Rungta) have stronger presence")

add_body_text(slide,
    CONTENT_LEFT, Emu(900000),
    CONTENT_WIDTH, Emu(5200000),
    paras)

print("  Slide 2: Dealer Growth Plan ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: Helix TMT Strategic Launch Plan
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Helix TMT: Strategic Launch Plan")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

# Left column: Market Positioning
left_paras = []
add_section_header(left_paras, "Market Positioning & Target Audience", spacing_before=Pt(0))
left_paras.append({'text': '', 'size': Pt(4)})
left_paras.append({
    'runs': [
        {'text': 'The "Gap": ', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Target the ', 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Economy segment', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': ' where JSW One is perceived as too expensive.', 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})
left_paras.append({
    'runs': [
        {'text': 'Niche: ', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Focus on ', 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Rural Markets', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': ' where primary steel demand is low and customers prioritize "physically and chemically approved" quality at a competitive price.', 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})

add_body_text(slide,
    Emu(332220), Emu(900000),
    Emu(5500000), Emu(5200000),
    left_paras)

# Right column: Pricing Strategy
right_paras = []
add_section_header(right_paras, "Pricing Strategy", spacing_before=Pt(0))
right_paras.append({'text': '', 'size': Pt(4)})
right_paras.append({
    'runs': [
        {'text': 'Benchmark: ', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Align directly with secondary brands like ', 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Jindal, Rapid, and Rungta.', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})
right_paras.append({
    'runs': [
        {'text': 'Formula: ', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
        {'text': 'Position at equivalent or minimal premium pricing to ensure the price reflects "Best Product at a Good Competitive Price" without hitting premium JSW levels.', 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})

add_body_text(slide,
    Emu(6200000), Emu(900000),
    Emu(5600000), Emu(5200000),
    right_paras)

print("  Slide 3: Helix Launch Strategy ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: Helix Channel Selection Criteria (Table)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Helix Channel Selection Criteria")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

# Create table
table_data = [
    ["Scenario", "Action Plan"],
    ["High-Performing JSW Dealers", "Onboard for Helix only if they have high financial capacity and untapped local market potential."],
    ["Limited Capacity Dealers", "Appoint new/alternative channel partners to ensure Helix gets the dedicated focus it needs."],
    ["Target Counters", "Focus on retailers already moving high volumes of Jindal/Rapid/Rungta."],
]

rows = len(table_data)
cols = 2
table_shape = slide.shapes.add_table(rows, cols, Emu(332220), Emu(1100000), Emu(11524664), Emu(2400000))
table = table_shape.table

# Set column widths
table.columns[0].width = Emu(3800000)
table.columns[1].width = Emu(7724664)

for ri, row_data in enumerate(table_data):
    for ci, cell_text in enumerate(row_data):
        cell = table.cell(ri, ci)
        if ri == 0:
            # Header row
            style_cell(cell, cell_text, bold=True, color=TABLE_HEADER_FG, bg=TABLE_HEADER_BG, font_size=Pt(12))
        else:
            bg = TABLE_ALT_BG if ri % 2 == 0 else None
            bold = (ci == 0)  # Bold the scenario name
            style_cell(cell, cell_text, bold=bold, color=BODY_COLOR, bg=bg, font_size=Pt(11))

print("  Slide 4: Channel Selection Table ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: Brand Differentiation (merged: Stocking + Brand Diff)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Brand Differentiation: JSW One vs One Helix")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

# Stocking section
left_paras = []
add_section_header(left_paras, "Stocking Strategy", spacing_before=Pt(0))
left_paras.append({'text': '', 'size': Pt(4)})
left_paras.append({
    'text': 'Currently no requirement to maintain stock at yard as supplies are consistent from plant. Bringing material to godown increases input cost and makes us less competitive. Will re-evaluate if required in future.',
    'size': BODY_SIZE,
    'spacing_before': Pt(8)
})

add_body_text(slide,
    Emu(332220), Emu(900000),
    Emu(5500000), Emu(5200000),
    left_paras)

# Brand differentiation
right_paras = []
add_section_header(right_paras, "Brand Positioning", spacing_before=Pt(0))
right_paras.append({'text': '', 'size': Pt(4)})

right_paras.append({
    'runs': [
        {'text': 'JSW One', 'bold': True, 'size': Pt(13), 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(8)
})
add_bullet(right_paras, "Premium positioning, priced higher")
add_bullet(right_paras, "Full services, quality product and marketing support")

right_paras.append({'text': '', 'size': Pt(6)})

right_paras.append({
    'runs': [
        {'text': 'One Helix', 'bold': True, 'size': Pt(13), 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(8)
})
add_bullet(right_paras, "Value-focused, competitive pricing")
add_bullet(right_paras, "Good services, wider channel reach")
add_bullet(right_paras, "Targets cost-sensitive customers")

add_body_text(slide,
    Emu(6200000), Emu(900000),
    Emu(5600000), Emu(5200000),
    right_paras)

print("  Slide 5: Brand Differentiation ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: Industry Best Practices
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Industry Best Practices: Lead & Secondary Brands")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

# Left: Primary brands
left_paras = []
add_section_header(left_paras, "2 Primary Brands: TATA, Jindal Panther", spacing_before=Pt(0))
left_paras.append({'text': '', 'size': Pt(4)})
left_paras.append({
    'runs': [
        {'text': 'Good Practices/Policies:', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})
add_bullet(left_paras, "Govt approvals (PWD)")
add_bullet(left_paras, "Influencer loyalty (Architects)")
add_bullet(left_paras, "Strong brand pull")

add_body_text(slide,
    Emu(332220), Emu(900000),
    Emu(5500000), Emu(5200000),
    left_paras)

# Right: Secondary brands
right_paras = []
add_section_header(right_paras, "2 Secondary Brands: Jindal, Rapid/Rungta", spacing_before=Pt(0))
right_paras.append({'text': '', 'size': Pt(4)})
right_paras.append({
    'runs': [
        {'text': 'Good Practices/Policies:', 'bold': True, 'size': BODY_SIZE, 'color': BODY_COLOR},
    ],
    'spacing_before': Pt(8)
})
add_bullet(right_paras, 'Strong "Sauda" booking culture')
add_bullet(right_paras, "Transparent schemes")
add_bullet(right_paras, "Pure price play")
add_bullet(right_paras, "High availability in rural belts / small supplies")

add_body_text(slide,
    Emu(6200000), Emu(900000),
    Emu(5600000), Emu(5200000),
    right_paras)

print("  Slide 6: Industry Best Practices ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: Influencer & Channel Meets Plan
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Influencer & Channel Meets Plan (FY26-27)")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

paras = []
add_section_header(paras, "District-wise Meet Identification & Frequency", spacing_before=Pt(0))
paras.append({'text': '', 'size': Pt(4)})

add_bullet(paras, "Nuh and Palwal: Very low sale of branded steel — plan One Helix for these markets. Mason meets where required.")
add_bullet(paras, "Influencer meets: On demand only, should not be linked to target or incentive.")
add_bullet(paras, "NCR Focus: Architects, engineers and builders are decision makers — plan effective activities for them.")
add_bullet(paras, "Requirement: Two TSEs for Gautam Budh Nagar and Hapur; One DSR for Project sales.")

paras.append({'text': '', 'size': Pt(8)})
add_section_header(paras, "Recommended Meet Frequency (North Haryana)")
paras.append({'text': '', 'size': Pt(4)})

# Table for frequency
meets_table_data = [
    ["Meet Type", "Frequency", "Level"],
    ["Mason Meets", "Every 2 months", "At counters"],
    ["Contractor Meets", "Every 2 months", "District level"],
    ["Architect Meetings", "Half-yearly / Yearly", "Regional"],
]

add_body_text(slide,
    CONTENT_LEFT, Emu(900000),
    CONTENT_WIDTH, Emu(3000000),
    paras)

# Add a small table for the meets frequency
table_shape = slide.shapes.add_table(4, 3, Emu(332220), Emu(4200000), Emu(8000000), Emu(1600000))
table = table_shape.table

table.columns[0].width = Emu(2666000)
table.columns[1].width = Emu(2666000)
table.columns[2].width = Emu(2666000)

for ri, row_data in enumerate(meets_table_data):
    for ci, cell_text in enumerate(row_data):
        cell = table.cell(ri, ci)
        if ri == 0:
            style_cell(cell, cell_text, bold=True, color=TABLE_HEADER_FG, bg=TABLE_HEADER_BG, font_size=Pt(11))
        else:
            bg = TABLE_ALT_BG if ri % 2 == 0 else None
            style_cell(cell, cell_text, color=BODY_COLOR, bg=bg, font_size=Pt(11))

print("  Slide 7: Meets Plan ✓")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: Distributor Key Performance Improvements
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(target_layout)
add_title(slide, "Distributor Key Performance Improvements (FY26-27)")
add_slide_number(slide, slide_num_counter)
slide_num_counter += 1

paras = []
add_section_header(paras, "Three Key Performance Points for Distributors", spacing_before=Pt(0))
paras.append({'text': '', 'size': Pt(6)})

# Point 1
paras.append({
    'runs': [
        {'text': '1', 'size': Pt(16), 'bold': True, 'color': SECTION_COLOR},
        {'text': '  Expand Market Reach & Sales Volume', 'size': Pt(14), 'bold': True, 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(16)
})
paras.append({
    'text': 'Increase retailer coverage and actively engage with contractors and construction projects for higher demand.',
    'size': BODY_SIZE,
    'spacing_before': Pt(6)
})

paras.append({'text': '', 'size': Pt(6)})

# Point 2
paras.append({
    'runs': [
        {'text': '2', 'size': Pt(16), 'bold': True, 'color': SECTION_COLOR},
        {'text': '  Appoint Dedicated Branch Manager', 'size': Pt(14), 'bold': True, 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(16)
})
paras.append({
    'text': 'Ensure focused leadership at branch level to drive sales execution and team management.',
    'size': BODY_SIZE,
    'spacing_before': Pt(6)
})

paras.append({'text': '', 'size': Pt(6)})

# Point 3
paras.append({
    'runs': [
        {'text': '3', 'size': Pt(16), 'bold': True, 'color': SECTION_COLOR},
        {'text': '  Improve Joint Visits with JSW One Team', 'size': Pt(14), 'bold': True, 'color': SECTION_COLOR},
    ],
    'spacing_before': Pt(16)
})
paras.append({
    'text': 'Strengthen collaboration with channel partners, key influencers, and architects through regular joint visits.',
    'size': BODY_SIZE,
    'spacing_before': Pt(6)
})

add_body_text(slide,
    CONTENT_LEFT, Emu(900000),
    CONTENT_WIDTH, Emu(5200000),
    paras)

print("  Slide 8: Distributor KPIs ✓")

# ── Save ────────────────────────────────────────────────────────────────
output_path = r'D:\RandomTestsClaude\Haryana_NewSlides.pptx'
prs.save(output_path)
print(f"\nSaved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
